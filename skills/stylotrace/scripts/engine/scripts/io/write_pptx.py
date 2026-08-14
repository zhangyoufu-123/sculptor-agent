#!/usr/bin/env python3
"""markdown（大纲式）→ pptx。用法: write_pptx.py <in.md> <out.pptx>

约定：
  # 标题        → 封面页（第一个标题）
  ## 小节标题   → 每页标题；## 之下的段落/列表项 → 该页要点（每页最多 8 条）
  纯文本段落    → 并入当前页要点
  无 ## 只有正文 → 按自然段分页
"""
import sys

from pptx import Presentation


def main(in_md, out_pptx):
    prs = Presentation()
    lines = open(in_md, encoding="utf-8").read().splitlines()
    current = None
    body = []
    cover_done = False

    def flush():
        nonlocal current, body, cover_done
        if current is None and not body:
            return
        if not cover_done and current and not body:
            slide = prs.slides.add_slide(prs.slide_layouts[0])
            slide.shapes.title.text = current[:80]
            cover_done = True
            current = None
            return
        slide = prs.slides.add_slide(prs.slide_layouts[1])
        slide.shapes.title.text = (current or "（无标题）")[:80]
        if body:
            tf = slide.placeholders[1].text_frame
            for i, item in enumerate(body[:8]):
                p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
                p.text = item[:120]
        body = []
        current = None
        cover_done = True

    for raw in lines:
        line = raw.rstrip()
        if not line.strip():
            continue
        if line.startswith("# "):
            flush()
            current = line[2:].strip()
        elif line.startswith("## "):
            flush()
            current = line[3:].strip()
        elif line.startswith("##"):
            flush()
            current = line.lstrip("#").strip()
        else:
            item = line.strip()
            if item.startswith("- "):
                item = item[2:]
            body.append(item)
    flush()
    prs.save(out_pptx)


if __name__ == "__main__":
    if len(sys.argv) < 3:
        sys.exit("用法: write_pptx.py <in.md> <out.pptx>")
    main(sys.argv[1], sys.argv[2])
